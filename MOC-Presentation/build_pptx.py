"""Build Mari MOC presentation as .pptx with brand colours, logo, animations enabled."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

HERE = Path(__file__).parent
LOGO = str(HERE / "logo.png")

# Mari palette
NAVY = RGBColor(0x0B, 0x25, 0x45)
BLUE = RGBColor(0x1E, 0x40, 0xAF)
TEAL = RGBColor(0x08, 0x91, 0xB2)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xDC, 0x26, 0x26)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0F, 0x17, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, shadow=False, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if not shadow:
        shape.shadow.inherit = False
    return shape


def add_logo(slide, x, y, h):
    return slide.shapes.add_picture(LOGO, x, y, height=h)


def footer(slide, page, total):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), fill=NAVY)
    add_text(slide, Inches(0.3), SH - Inches(0.33), Inches(6), Inches(0.3),
             "Mari Energies  |  MOC Management System", size=10, color=WHITE)
    add_text(slide, SW - Inches(2), SH - Inches(0.33), Inches(1.7), Inches(0.3),
             f"{page} / {total}", size=10, color=WHITE, align=PP_ALIGN.RIGHT)


# ---------------- Slide 1 — Cover ----------------
TOTAL = 22

s = add_slide(NAVY)
# decorative band
add_rect(s, 0, Inches(2.0), SW, Inches(3.5), fill=RGBColor(0x12, 0x2E, 0x55))
add_logo(s, Inches(5.66), Inches(0.7), Inches(1.4))
add_text(s, 0, Inches(2.4), SW, Inches(0.9),
         "MOC Management System", size=44, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(3.3), SW, Inches(0.6),
         "Management of Change — MSP-HSE-08", size=22, color=GOLD,
         align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(4.1), SW, Inches(0.5),
         "A digital platform for safe, auditable change at Mari Energies",
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.16), Inches(5.2), Inches(3.0), Inches(0.55), fill=GOLD, radius=True)
add_text(s, Inches(5.16), Inches(5.22), Inches(3.0), Inches(0.5),
         "mari.proflowenergy.org", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------------- Slide 2 — Agenda ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Agenda", size=28, bold=True, color=WHITE)

agenda = [
    ("01", "The Problem", "Why we need a digital MOC system"),
    ("02", "MSP-HSE-08 Compliance", "How we map to the procedure"),
    ("03", "Solution Overview", "Platform at a glance"),
    ("04", "Core Features", "What it does end-to-end"),
    ("05", "Hierarchy & Workflow", "Approval chains & delegation"),
    ("06", "Forms & Reports", "Annexure H, Minute Sheet, Audit"),
    ("07", "Security & Tech", "Stack, deployment, controls"),
    ("08", "Roadmap & Q&A", "What's next"),
]
for i, (num, title, sub) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.4 + row * 1.35)
    add_rect(s, x, y, Inches(5.9), Inches(1.15), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(0.15), Inches(1.15), fill=GOLD)
    add_text(s, x + Inches(0.35), y + Inches(0.1), Inches(1.0), Inches(0.5),
             num, size=24, bold=True, color=GOLD)
    add_text(s, x + Inches(1.4), y + Inches(0.15), Inches(4.4), Inches(0.4),
             title, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(1.4), y + Inches(0.6), Inches(4.4), Inches(0.5),
             sub, size=11, color=SLATE)
footer(s, 2, TOTAL)

# ---------------- Slide 3 — Problem ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "The Problem with Paper-based MOC", size=24, bold=True, color=WHITE)

problems = [
    ("Lost paperwork", "Forms get misfiled. Approvals stall for weeks.", RED),
    ("No audit trail", "Who approved what, when? Hard to reconstruct.", GOLD),
    ("Manual hierarchy", "Routing forms by hand — error prone & slow.", BLUE),
    ("No analytics", "No KPIs, no closure timing, no trend visibility.", TEAL),
    ("Compliance risk", "MSP-HSE-08 § audit findings on traceability.", RED),
    ("Knowledge silos", "Field engineers can't see status of their MOC.", GOLD),
]
for i, (t, d, c) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(1.4 + row * 2.8)
    add_rect(s, x, y, Inches(4.0), Inches(2.55), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(4.0), Inches(0.5), fill=c)
    add_text(s, x + Inches(0.2), y + Inches(0.05), Inches(3.7), Inches(0.45),
             t, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.7), Inches(1.7),
             d, size=12, color=SLATE)
footer(s, 3, TOTAL)

# ---------------- Slide 4 — MSP-HSE-08 Compliance ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Built on MSP-HSE-08 Procedure", size=24, bold=True, color=WHITE)

mappings = [
    ("§ 3.4.1", "Initiation", "Originator raises MOC with description, justification, classification"),
    ("§ 3.4.2", "Risk Screening", "Annexure C/D1/D2 — Initial Risk, Classification, Risk Level"),
    ("§ 3.4.3", "Hierarchy Routing", "Approval chains for Minor / Major-Capital / Major-NonCapital / Field Pkg"),
    ("§ 3.4.4", "Detailed Engineering", "Stage-2 functional reviews & Work Pack"),
    ("§ 3.4.5", "Execution & Closure", "Implementation tracking, residual risk, close-out"),
    ("§ 3.5", "Records & Audit", "Full digital trail, exportable Minute Sheet & Audit Report"),
]
y0 = Inches(1.4)
for i, (sec, title, desc) in enumerate(mappings):
    y = y0 + Inches(i * 0.85)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.78), fill=LIGHT, radius=True)
    add_rect(s, Inches(0.5), y, Inches(1.4), Inches(0.78), fill=NAVY, radius=True)
    add_text(s, Inches(0.5), y, Inches(1.4), Inches(0.78),
             sec, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.1), y + Inches(0.08), Inches(2.8), Inches(0.4),
             title, size=14, bold=True, color=NAVY)
    add_text(s, Inches(2.1), y + Inches(0.42), Inches(10.5), Inches(0.4),
             desc, size=11, color=SLATE)
footer(s, 4, TOTAL)

# ---------------- Slide 5 — Solution Overview ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Solution Overview", size=24, bold=True, color=WHITE)

# Mock UI: sidebar + dashboard tiles
add_rect(s, Inches(0.5), Inches(1.4), Inches(2.4), Inches(5.4), fill=NAVY, radius=True)
add_logo(s, Inches(0.95), Inches(1.6), Inches(0.5))
add_text(s, Inches(0.6), Inches(2.3), Inches(2.2), Inches(0.4),
         "Mari MOC", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
nav_items = ["Dashboard", "MOC Register", "My Tasks", "Audit", "Users", "Settings"]
for i, item in enumerate(nav_items):
    y = Inches(2.9 + i * 0.42)
    fill = BLUE if i == 0 else NAVY
    add_rect(s, Inches(0.7), y, Inches(2.0), Inches(0.34), fill=fill, radius=True)
    add_text(s, Inches(0.85), y, Inches(1.9), Inches(0.34),
             item, size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# KPI tiles
kpis = [("47", "Total MOCs", BLUE), ("12", "In Review", GOLD),
        ("8", "Approved", GREEN), ("3", "Rejected", RED)]
for i, (val, label, c) in enumerate(kpis):
    x = Inches(3.1 + i * 2.45)
    add_rect(s, x, Inches(1.4), Inches(2.3), Inches(1.4), fill=LIGHT, radius=True)
    add_rect(s, x, Inches(1.4), Inches(0.12), Inches(1.4), fill=c)
    add_text(s, x + Inches(0.2), Inches(1.55), Inches(2.0), Inches(0.7),
             val, size=32, bold=True, color=c)
    add_text(s, x + Inches(0.2), Inches(2.25), Inches(2.0), Inches(0.4),
             label, size=11, color=SLATE)

# Mock table
add_rect(s, Inches(3.1), Inches(2.95), Inches(9.7), Inches(3.85), fill=LIGHT, radius=True)
add_text(s, Inches(3.3), Inches(3.05), Inches(9), Inches(0.4),
         "Recent MOC Register", size=14, bold=True, color=NAVY)
headers = ["MOC #", "Title", "Class", "Stage", "Status"]
hx = [3.3, 4.5, 7.2, 8.4, 10.5]
for i, h in enumerate(headers):
    add_text(s, Inches(hx[i]), Inches(3.5), Inches(2), Inches(0.3),
             h, size=10, bold=True, color=SLATE)
rows = [
    ("MOC-2026-014", "Replace PSV-301", "Major", "Stage 3", "In Review", GOLD),
    ("MOC-2026-013", "Add temp bypass", "Minor", "Stage 4", "Approved", GREEN),
    ("MOC-2026-012", "Pump impeller", "Minor", "Stage 2", "In Review", GOLD),
    ("MOC-2026-011", "DCS logic edit", "Major", "Stage 1", "Draft", SLATE),
    ("MOC-2026-010", "Insulation", "Field", "Closed", "Closed", BLUE),
]
for i, row in enumerate(rows):
    y = Inches(3.9 + i * 0.5)
    if i % 2 == 0:
        add_rect(s, Inches(3.2), y, Inches(9.5), Inches(0.45), fill=WHITE, radius=True)
    for j, val in enumerate(row[:5]):
        add_text(s, Inches(hx[j]), y + Inches(0.05), Inches(2), Inches(0.4),
                 val, size=10, color=DARK)
    # Status pill
    add_rect(s, Inches(10.5), y + Inches(0.08), Inches(1.0), Inches(0.3),
             fill=row[5], radius=True)
    add_text(s, Inches(10.5), y + Inches(0.08), Inches(1.0), Inches(0.3),
             row[4], size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5, TOTAL)

# ---------------- Slide 6 — Core Features ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "12 Core Features", size=24, bold=True, color=WHITE)

features = [
    ("Initiation", "Web form for any engineer", BLUE),
    ("Classification", "Annexure C/D1/D2 wizards", TEAL),
    ("Risk Scoring", "Initial & residual risk matrices", GOLD),
    ("Hierarchy Routing", "Auto-builds approval chain", GREEN),
    ("Delegation", "Assign to qualified subordinate", BLUE),
    ("Stage Forms", "Stage 1-4 with role-locked fields", TEAL),
    ("Annexure H", "Auto-generated 3-page CRF", GOLD),
    ("Minute Sheet", "PDF audit trail per MOC", GREEN),
    ("Dashboards", "KPIs, timing, breakdowns", BLUE),
    ("Audit Report", "Director-level analytics", TEAL),
    ("Notifications", "Inbox + email on action", GOLD),
    ("Mobile Ready", "Responsive across devices", GREEN),
]
for i, (t, d, c) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.2)
    y = Inches(1.4 + row * 1.85)
    add_rect(s, x, y, Inches(3.0), Inches(1.65), fill=LIGHT, radius=True)
    add_rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
             fill=c, radius=True)
    add_text(s, x + Inches(0.85), y + Inches(0.2), Inches(2.0), Inches(0.45),
             t, size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.85), Inches(2.7), Inches(0.7),
             d, size=11, color=SLATE)
footer(s, 6, TOTAL)

# ---------------- Slide 7 — 4-Stage Lifecycle ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "MOC Lifecycle — 4 Stages", size=24, bold=True, color=WHITE)

stages = [
    ("STAGE 1", "Initiate", "Originator fills request, attaches drawings, classifies", BLUE),
    ("STAGE 2", "Review & Approve", "Hierarchy chain reviews, comments, approves or rejects", TEAL),
    ("STAGE 3", "Execute", "Approved MOC released for implementation under Work Pack", GOLD),
    ("STAGE 4", "Close-out", "Residual risk, lessons learned, archive in audit log", GREEN),
]
for i, (n, t, d, c) in enumerate(stages):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.6)
    add_rect(s, x, y, Inches(3.0), Inches(4.5), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(3.0), Inches(0.7), fill=c)
    add_text(s, x, y, Inches(3.0), Inches(0.7),
             n, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.9), Inches(2.6), Inches(0.5),
             t, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.6), Inches(2.6), Inches(2.5),
             d, size=12, color=SLATE)
    if i < 3:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + Inches(3.0), y + Inches(2.0),
                                 Inches(0.2), Inches(0.5))
        arr.fill.solid(); arr.fill.fore_color.rgb = GOLD
        arr.line.fill.background()
footer(s, 7, TOTAL)

# ---------------- Slide 8 — Hierarchy Pyramid ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Approval Hierarchy", size=24, bold=True, color=WHITE)

# Pyramid (5 levels)
levels = [
    ("HO Director Operations / Head EDP", NAVY, 4.0),
    ("Manager Process Ops / HSE / Engineering", BLUE, 5.5),
    ("Manager MAI / MOC Interface", TEAL, 7.0),
    ("Manager Production", GREEN, 8.5),
    ("Field In-Charge → Originator (JRE)", GOLD, 10.0),
]
y0 = 1.5
for i, (label, c, w) in enumerate(levels):
    x = (13.333 - w) / 2
    y = y0 + i * 0.85
    add_rect(s, Inches(x), Inches(y), Inches(w), Inches(0.75), fill=c, radius=True)
    add_text(s, Inches(x), Inches(y), Inches(w), Inches(0.75),
             label, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
         "Each MOC is auto-routed up the chain. Approvals are forwarded; rejections return to originator with comments.",
         size=12, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, 8, TOTAL)

# ---------------- Slide 9 — Minor vs Major chains ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Minor vs Major Chain", size=24, bold=True, color=WHITE)

minor_chain = ["Originator (JRE)", "Field In-Charge", "Mgr Production",
               "MOC Interface", "Mgr MAI", "Eng Manager",
               "Mgr HSE", "Mgr Process Ops", "Director Ops"]
major_chain = ["Originator (JRE)", "Field In-Charge", "Mgr Production",
               "MOC Interface", "Mgr MAI", "Eng Manager",
               "Mgr HSE", "Mgr Process Ops", "Director HSE",
               "Director Ops / Head EDP"]

for col, (title, chain, c) in enumerate([
    ("MINOR (9 steps)", minor_chain, GREEN),
    ("MAJOR (10 steps)", major_chain, GOLD),
]):
    x = Inches(0.5 + col * 6.4)
    add_rect(s, x, Inches(1.4), Inches(6.2), Inches(5.5), fill=LIGHT, radius=True)
    add_rect(s, x, Inches(1.4), Inches(6.2), Inches(0.6), fill=c)
    add_text(s, x, Inches(1.4), Inches(6.2), Inches(0.6),
             title, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, step in enumerate(chain):
        y = Inches(2.15 + i * 0.45)
        add_rect(s, x + Inches(0.3), y, Inches(0.45), Inches(0.35),
                 fill=NAVY, radius=True)
        add_text(s, x + Inches(0.3), y, Inches(0.45), Inches(0.35),
                 str(i + 1), size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), y, Inches(5.2), Inches(0.35),
                 step, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 9, TOTAL)

# ---------------- Slide 10 — Delegation ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Delegation", size=24, bold=True, color=WHITE)

add_text(s, Inches(0.6), Inches(1.3), Inches(7), Inches(0.5),
         "When the assigned approver is unavailable", size=16, bold=True, color=NAVY)

bullets = [
    "Approver opens the MOC and clicks Delegate to Subordinate",
    "System lists qualified direct reports (same dept or trained)",
    "Selected subordinate receives the task in their inbox",
    "Delegation is recorded permanently in the audit trail",
    "Original chain continues from there — no skipping levels",
]
for i, b in enumerate(bullets):
    y = Inches(2.0 + i * 0.55)
    add_rect(s, Inches(0.6), y + Inches(0.1), Inches(0.2), Inches(0.2),
             fill=GOLD, radius=True)
    add_text(s, Inches(1.0), y, Inches(6.5), Inches(0.5),
             b, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# Mock delegate panel
x0 = Inches(8.0)
add_rect(s, x0, Inches(1.5), Inches(4.8), Inches(4.5), fill=LIGHT, radius=True)
add_rect(s, x0, Inches(1.5), Inches(4.8), Inches(0.6), fill=NAVY, radius=True)
add_text(s, x0, Inches(1.5), Inches(4.8), Inches(0.6),
         "Delegate Task", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x0 + Inches(0.3), Inches(2.3), Inches(4.4), Inches(0.4),
         "Select subordinate:", size=11, color=SLATE)
subs = [("Asad — Field In-Charge", True),
        ("Tariq — Mgr Production", False),
        ("Bilal — Mgr MAI", False)]
for i, (name, sel) in enumerate(subs):
    y = Inches(2.75 + i * 0.55)
    add_rect(s, x0 + Inches(0.3), y, Inches(4.2), Inches(0.45),
             fill=WHITE if not sel else BLUE, radius=True)
    add_text(s, x0 + Inches(0.5), y, Inches(4.0), Inches(0.45),
             name, size=11, bold=sel, color=WHITE if sel else DARK,
             anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, x0 + Inches(0.3), Inches(5.2), Inches(4.2), Inches(0.6),
         fill=GOLD, radius=True)
add_text(s, x0 + Inches(0.3), Inches(5.2), Inches(4.2), Inches(0.6),
         "Confirm Delegation", size=12, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10, TOTAL)

# ---------------- Slide 11 — Reject Flow ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Reject / Revision Flow", size=24, bold=True, color=WHITE)

flow = [
    ("Approver Reviews", BLUE),
    ("Marks 'Revision Required'", GOLD),
    ("Comments captured", GOLD),
    ("Returns to Originator", RED),
    ("Originator edits & resubmits", GREEN),
    ("Chain restarts from Step 1", TEAL),
]
for i, (t, c) in enumerate(flow):
    x = Inches(0.5 + (i % 3) * 4.3)
    y = Inches(1.7 + (i // 3) * 2.4)
    add_rect(s, x, y, Inches(4.0), Inches(2.0), fill=LIGHT, radius=True)
    add_rect(s, x + Inches(1.5), y - Inches(0.3), Inches(1.0), Inches(1.0),
             fill=c, radius=True)
    add_text(s, x + Inches(1.5), y - Inches(0.3), Inches(1.0), Inches(1.0),
             str(i + 1), size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + Inches(0.9), Inches(4.0), Inches(0.9),
             t, size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11, TOTAL)

# ---------------- Slide 12 — Stage Forms ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Stage Forms — Role-locked Fields", size=24, bold=True, color=WHITE)

stage_forms = [
    ("Stage 1", "Description, Justification, Drawings, Classification (Annex C)", BLUE),
    ("Stage 2", "Risk Screening (D1), Risk Level (D2), Functional Reviews", TEAL),
    ("Stage 3", "Detailed Engineering, Work Pack, Execution Plan", GOLD),
    ("Stage 4", "Implementation Sign-off, Residual Risk, Lessons Learned, Close-out", GREEN),
]
for i, (n, d, c) in enumerate(stage_forms):
    y = Inches(1.4 + i * 1.35)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.2), fill=LIGHT, radius=True)
    add_rect(s, Inches(0.5), y, Inches(1.8), Inches(1.2), fill=c, radius=True)
    add_text(s, Inches(0.5), y, Inches(1.8), Inches(1.2),
             n, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), y + Inches(0.2), Inches(10), Inches(0.4),
             "Required fields:", size=11, bold=True, color=SLATE)
    add_text(s, Inches(2.5), y + Inches(0.6), Inches(10), Inches(0.6),
             d, size=12, color=DARK)
footer(s, 12, TOTAL)

# ---------------- Slide 13 — Annexure H / Minute Sheet / Audit ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Reports & Outputs", size=24, bold=True, color=WHITE)

reports = [
    ("Annexure H",
     "Official Mari Change Request Form\n\n• Page 1: CRF\n• Page 2: Stage-2 Engg\n• Page 3: Work Pack\n\nAuto-populated from MOC data — preview & print",
     BLUE),
    ("Minute Sheet",
     "PDF audit trail\n\n• All approvers + timestamps\n• Comments at each step\n• Delegations captured\n• Status transitions\n\nDownloaded from MOC detail",
     GOLD),
    ("Audit Report",
     "Director-level analytics\n\n• KPIs: total, in-review, closed, expired\n• Avg closure time\n• Breakdowns by class, dept, risk\n• Activity timeline",
     GREEN),
]
for i, (t, d, c) in enumerate(reports):
    x = Inches(0.5 + i * 4.3)
    add_rect(s, x, Inches(1.5), Inches(4.0), Inches(5.0), fill=LIGHT, radius=True)
    add_rect(s, x, Inches(1.5), Inches(4.0), Inches(0.7), fill=c, radius=True)
    add_text(s, x, Inches(1.5), Inches(4.0), Inches(0.7),
             t, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), Inches(2.4), Inches(3.7), Inches(4.0),
             d, size=12, color=DARK)
footer(s, 13, TOTAL)

# ---------------- Slide 14 — Director Dashboard KPIs ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Director Dashboard — KPIs", size=24, bold=True, color=WHITE)

kpis2 = [
    ("47", "Total MOCs", BLUE),
    ("12", "In Review", GOLD),
    ("28", "Closed", GREEN),
    ("3", "Rejected", RED),
    ("4", "Expired", RGBColor(0xEA, 0x58, 0x0C)),
    ("18d", "Avg Closure", TEAL),
    ("92%", "Compliance", GREEN),
    ("6", "Pending Mine", BLUE),
]
for i, (val, label, c) in enumerate(kpis2):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.2)
    y = Inches(1.4 + row * 1.7)
    add_rect(s, x, y, Inches(3.0), Inches(1.5), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(0.15), Inches(1.5), fill=c)
    add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2.6), Inches(0.7),
             val, size=32, bold=True, color=c)
    add_text(s, x + Inches(0.3), y + Inches(0.95), Inches(2.6), Inches(0.4),
             label, size=12, color=SLATE)

# Bar chart row
add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.85), fill=LIGHT, radius=True)
add_text(s, Inches(0.7), Inches(5.1), Inches(10), Inches(0.4),
         "MOCs by Department", size=13, bold=True, color=NAVY)
depts = [("OPS", 0.85, BLUE), ("ENG", 0.65, TEAL),
         ("HSE", 0.45, GREEN), ("MAI", 0.55, GOLD), ("EDP", 0.30, RED)]
for i, (d, w, c) in enumerate(depts):
    y = Inches(5.55 + i * 0.25)
    add_text(s, Inches(0.7), y, Inches(0.8), Inches(0.22),
             d, size=10, bold=True, color=DARK)
    add_rect(s, Inches(1.6), y + Inches(0.03), Inches(10 * w), Inches(0.18),
             fill=c, radius=True)
footer(s, 14, TOTAL)

# ---------------- Slide 15 — Roles & Permissions ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Roles & Permissions", size=24, bold=True, color=WHITE)

roles = [
    ("Engineer", "Raise MOC, edit own drafts, view own MOCs, resubmit", BLUE),
    ("Approver", "Review assigned MOCs, approve / reject / delegate", TEAL),
    ("Department Head", "All approver rights + dept analytics", GOLD),
    ("Director", "All MOCs visible, audit report, executive KPIs", GREEN),
    ("Admin", "User mgmt, hierarchy config, system settings", RED),
]
for i, (r, p, c) in enumerate(roles):
    y = Inches(1.4 + i * 1.05)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.95), fill=LIGHT, radius=True)
    add_rect(s, Inches(0.5), y, Inches(2.5), Inches(0.95), fill=c, radius=True)
    add_text(s, Inches(0.5), y, Inches(2.5), Inches(0.95),
             r, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.2), y, Inches(9.4), Inches(0.95),
             p, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 15, TOTAL)

# ---------------- Slide 16 — Security ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Security & Compliance", size=24, bold=True, color=WHITE)

sec = [
    ("HTTPS Everywhere", "TLS 1.2+, Let's Encrypt auto-renew", GREEN),
    ("JWT Auth", "Signed tokens, 7-day expiry, rotated secret", BLUE),
    ("Bcrypt Passwords", "Salted hashes, 12 rounds, never plaintext", TEAL),
    ("Role-based Access", "Field-level locks, every endpoint authorized", GOLD),
    ("Audit Trail", "Immutable log: who, what, when, why", RED),
    ("Daily Backups", "Automated DB dumps with off-site copy", GREEN),
]
for i, (t, d, c) in enumerate(sec):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(1.5 + row * 2.7)
    add_rect(s, x, y, Inches(4.0), Inches(2.4), fill=LIGHT, radius=True)
    add_rect(s, x + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7),
             fill=c, radius=True)
    add_text(s, x + Inches(1.2), y + Inches(0.35), Inches(2.6), Inches(0.7),
             t, size=14, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), y + Inches(1.2), Inches(3.5), Inches(1.1),
             d, size=11, color=SLATE)
footer(s, 16, TOTAL)

# ---------------- Slide 17 — Tech Stack ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Technology Stack", size=24, bold=True, color=WHITE)

stack = [
    ("Frontend", "Vanilla JS SPA\nNo build step\nResponsive CSS", BLUE),
    ("Backend", "Node.js 20 + Express\nSequelize ORM\nJWT auth", TEAL),
    ("Database", "MariaDB 10.5+\nFully relational\n8 core tables", GOLD),
    ("Hosting", "DirectAdmin / VPS\nNginx + pm2\nLet's Encrypt TLS", GREEN),
]
for i, (t, d, c) in enumerate(stack):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(1.5), Inches(3.0), Inches(4.5), fill=LIGHT, radius=True)
    add_rect(s, x, Inches(1.5), Inches(3.0), Inches(0.8), fill=c, radius=True)
    add_text(s, x, Inches(1.5), Inches(3.0), Inches(0.8),
             t, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), Inches(2.6), Inches(2.7), Inches(3.0),
             d, size=13, color=DARK, align=PP_ALIGN.CENTER)
footer(s, 17, TOTAL)

# ---------------- Slide 18 — Deployment ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Deployment", size=24, bold=True, color=WHITE)

deploy = [
    ("1", "Provision Server", "VPS or DirectAdmin shared hosting"),
    ("2", "Restore Database", "Import db_dump.sql via mysql / phpMyAdmin"),
    ("3", "Configure .env", "DB creds, JWT_SECRET (rotate!), PUBLIC_URL"),
    ("4", "Install + Start", "npm install --production then pm2 / DirectAdmin"),
    ("5", "Frontend Static", "Copy frontend/ to public_html, point API_BASE"),
    ("6", "TLS Certificate", "certbot --nginx -d mari.proflowenergy.org"),
    ("7", "Smoke Test", "curl /api/mocs → 401 means API live"),
    ("8", "Daily Backups", "mysqldump cron → off-site"),
]
for i, (n, t, d) in enumerate(deploy):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.4 + row * 1.35)
    add_rect(s, x, y, Inches(6.2), Inches(1.2), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(0.9), Inches(1.2), fill=NAVY, radius=True)
    add_text(s, x, y, Inches(0.9), Inches(1.2),
             n, size=28, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.1), y + Inches(0.18), Inches(5.0), Inches(0.4),
             t, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(1.1), y + Inches(0.6), Inches(5.0), Inches(0.55),
             d, size=11, color=SLATE)
footer(s, 18, TOTAL)

# ---------------- Slide 19 — Roadmap ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Roadmap", size=24, bold=True, color=WHITE)

roadmap = [
    ("Q2 2026", "Email + SMS notifications", GREEN),
    ("Q2 2026", "Mobile native app (iOS/Android)", BLUE),
    ("Q3 2026", "SSO / Active Directory integration", TEAL),
    ("Q3 2026", "Advanced analytics + ML risk scoring", GOLD),
    ("Q4 2026", "API for SAP / Maximo integration", RED),
    ("Q4 2026", "Multi-site rollout & i18n", NAVY),
]
for i, (q, t, c) in enumerate(roadmap):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(1.6 + row * 2.7)
    add_rect(s, x, y, Inches(4.0), Inches(2.4), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(4.0), Inches(0.6), fill=c, radius=True)
    add_text(s, x, y, Inches(4.0), Inches(0.6),
             q, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), y + Inches(0.9), Inches(3.5), Inches(1.4),
             t, size=14, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 19, TOTAL)

# ---------------- Slide 20 — Anticipated Q&A ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Anticipated Questions", size=24, bold=True, color=WHITE)

qa = [
    ("How long to deploy?", "Bundle ready — under 1 day on existing infra."),
    ("What if internet is down?", "Runs on internal LAN; offline bundle ships with full source."),
    ("Can we customise the chain?", "Yes — hierarchy is data-driven, edit per dept."),
    ("How is data backed up?", "Daily mysqldump + off-site copy; restore in minutes."),
    ("Who owns the source code?", "Mari Energies — full IP transfer with offline bundle."),
    ("Can it integrate with SAP?", "Yes — REST API ready; on Q4 roadmap."),
]
for i, (q, a) in enumerate(qa):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.4 + row * 1.85)
    add_rect(s, x, y, Inches(6.2), Inches(1.7), fill=LIGHT, radius=True)
    add_rect(s, x, y, Inches(0.5), Inches(1.7), fill=GOLD, radius=True)
    add_text(s, x, y, Inches(0.5), Inches(1.7),
             "Q", size=22, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.7), y + Inches(0.15), Inches(5.4), Inches(0.5),
             q, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.7), y + Inches(0.7), Inches(5.4), Inches(0.95),
             a, size=11, color=SLATE)
footer(s, 20, TOTAL)

# ---------------- Slide 21 — Value / ROI ----------------
s = add_slide()
add_rect(s, 0, 0, SW, Inches(1.0), fill=NAVY)
add_logo(s, Inches(0.4), Inches(0.18), Inches(0.65))
add_text(s, Inches(1.3), Inches(0.2), Inches(10), Inches(0.6),
         "Value & ROI", size=24, bold=True, color=WHITE)

values = [
    ("60%", "Faster MOC closure", "From weeks to days", GREEN),
    ("100%", "Audit trail coverage", "Every action logged", BLUE),
    ("Zero", "Lost paperwork", "Digital from Day 1", GOLD),
    ("24/7", "Anywhere access", "Web + mobile", TEAL),
    ("MSP-HSE-08", "Full compliance", "By design", RED),
    ("∞", "Scalable", "All sites, all classes", NAVY),
]
for i, (val, t, d, c) in enumerate(values):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(1.5 + row * 2.7)
    add_rect(s, x, y, Inches(4.0), Inches(2.4), fill=LIGHT, radius=True)
    add_text(s, x, y + Inches(0.2), Inches(4.0), Inches(0.9),
             val, size=44, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.2), Inches(4.0), Inches(0.5),
             t, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.7), Inches(4.0), Inches(0.5),
             d, size=11, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, 21, TOTAL)

# ---------------- Slide 22 — Thank You ----------------
s = add_slide(NAVY)
add_rect(s, 0, Inches(2.5), SW, Inches(2.5), fill=RGBColor(0x12, 0x2E, 0x55))
add_logo(s, Inches(5.66), Inches(0.9), Inches(1.4))
add_text(s, 0, Inches(2.7), SW, Inches(0.9),
         "Thank You", size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(3.7), SW, Inches(0.5),
         "Questions & Discussion", size=20, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.5), SW, Inches(0.4),
         "mari.proflowenergy.org", size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.95), SW, Inches(0.4),
         "MSP-HSE-08  |  Mari Energies MOC Management System",
         size=12, color=LIGHT, align=PP_ALIGN.CENTER)

out = HERE / "MOC-Presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
